# ✅ Full Streamlit App: Graph generated strictly from Gemini's extracted CSV (from image)

import os
import pandas as pd
import requests
import streamlit as st
from io import StringIO, BytesIO
from PIL import Image, UnidentifiedImageError
import streamlit as st
from dotenv import load_dotenv
import google.generativeai as genai
from atlassian import Confluence
from bs4 import BeautifulSoup
from fpdf import FPDF
from docx import Document
from docx.shared import Inches
import tempfile
import seaborn as sns
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches
import streamlit.components.v1 as components

load_dotenv()

@st.cache_resource

def init_confluence():
    try:
        return Confluence(
            url=os.getenv('CONFLUENCE_BASE_URL'),
            username=os.getenv('CONFLUENCE_USER_EMAIL'),
            password=os.getenv('CONFLUENCE_API_KEY'),
            timeout=10
        )
    except Exception as e:
        st.error(f"Confluence initialization failed: {str(e)}")
        return None

def init_ai():
    genai.configure(api_key=os.getenv("GENAI_API_KEY"))
    return genai.GenerativeModel("models/gemini-1.5-flash")

def download_image_bytes(url: str, auth):
    response = requests.get(url, auth=auth)
    if response.status_code == 200:
        return response.content
    else:
        raise Exception(f"Failed to fetch image (Status: {response.status_code})")

def generate_pdf(image_bytes, summary):
    pdf = FPDF()
    pdf.add_page()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        img = Image.open(BytesIO(image_bytes))
        img.save(tmp_img, format="PNG")
        tmp_img_path = tmp_img.name
    pdf.image(tmp_img_path, x=10, w=180)
    pdf.ln(5)
    safe_summary = summary.encode('latin-1', 'replace').decode('latin-1')
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, safe_summary)
    pdf_output = pdf.output(dest='S').encode('latin-1')
    return BytesIO(pdf_output)

def generate_docx(image_bytes, summary):
    doc = Document()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        img = Image.open(BytesIO(image_bytes))
        img.save(tmp_img, format="PNG")
        tmp_img_path = tmp_img.name
    doc.add_picture(tmp_img_path, width=Inches(5.5))
    doc.add_paragraph(summary)
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def generate_txt(summary):
    output = BytesIO()
    output.write(summary.encode())
    output.seek(0)
    return output

def generate_md(image_bytes, summary):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        img = Image.open(BytesIO(image_bytes))
        img.save(tmp_img, format="PNG")
        image_path = tmp_img.name
    content = f"![Image]({image_path})\n\n{summary}"
    output = BytesIO()
    output.write(content.encode())
    output.seek(0)
    return output

def plot_grouped_bar(df):
    melted = df.melt(id_vars=[df.columns[0]], var_name="Group", value_name="Count")
    plt.figure(figsize=(10, 6))
    sns.barplot(data=melted, x=melted.columns[0], y="Count", hue="Group")
    plt.xticks(rotation=45)
    plt.title("Grouped Bar Chart")
    plt.tight_layout()
    return plt.gcf()

def plot_stacked_bar(df):
    df_plot = df.set_index(df.columns[0])
    plt.figure(figsize=(10, 6))
    df_plot.drop(columns="Total", errors="ignore").plot(kind='bar', stacked=True)
    plt.title("Stacked Bar Chart")
    plt.xticks(rotation=45)
    plt.ylabel("Count")
    plt.tight_layout()
    return plt.gcf()

def plot_line(df):
    df_plot = df.set_index(df.columns[0])
    plt.figure(figsize=(10, 6))
    df_plot.drop(columns="Total", errors="ignore").plot(marker='o')
    plt.title("Line Chart")
    plt.xticks(rotation=45)
    plt.ylabel("Count")
    plt.tight_layout()
    return plt.gcf()

def plot_pie(df):
    plt.figure(figsize=(7, 6))
    label_col = df.columns[0]
    
    if "Total" in df.columns:
        data = df["Total"]
    else:
        # Fallback: Sum across all numeric columns (except label)
        data = df.iloc[:, 1:].sum(axis=1)

    plt.pie(data, labels=df[label_col], autopct="%1.1f%%", startangle=140)
    plt.title("Pie Chart (Total Responses)")
    plt.tight_layout()
    return plt.gcf()


def get_image_bytes(fig, fmt):
    buf = BytesIO()
    fig.savefig(buf, format=fmt, bbox_inches="tight")
    buf.seek(0)
    return buf

def clean_ai_csv(raw_text):
    lines = raw_text.strip().splitlines()
    clean_lines = [
        line.strip() for line in lines
        if ',' in line and not line.strip().startswith("```") and not line.lower().startswith("here")
    ]
    # Remove duplicate headers or malformed lines
    header = clean_lines[0].split(",")
    cleaned_data = [clean_lines[0]]
    for line in clean_lines[1:]:
        if line.split(",")[0] != header[0]:  # skip repeat headers
            cleaned_data.append(line)
    return "\n".join(cleaned_data)

def generate_chart_pdf(image_bytes):
    pdf = FPDF()
    pdf.add_page()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        tmp_img.write(image_bytes.read())
        tmp_img.flush()
        pdf.image(tmp_img.name, x=10, w=180)
    output = BytesIO()
    pdf.output(output)
    output.seek(0)
    return output

def generate_chart_docx(image_bytes):
    doc = Document()
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        tmp_img.write(image_bytes.read())
        tmp_img.flush()
        doc.add_picture(tmp_img.name, width=Inches(5.5))
    output = BytesIO()
    doc.save(output)
    output.seek(0)
    return output

def generate_chart_pptx(image_bytes):
    prs = Presentation()
    blank_slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
        tmp_img.write(image_bytes.read())
        tmp_img.flush()
        slide_width = prs.slide_width
        top = Inches(1)
        left = Inches(0.5)
        pic = blank_slide.shapes.add_picture(tmp_img.name, left, top, width=slide_width - Inches(1))
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# === Streamlit app ===

st.set_page_config(page_title="Confluence Image Summarizer", layout="wide")
st.title("Confluence Page Image Summarizer")



confluence = init_confluence()
ai_model = init_ai()

if confluence:
    st.success("✅ Connected to Confluence")

    space_key = st.text_input("Enter your space key:")

    if space_key:
        try:
            pages = confluence.get_all_pages_from_space(space=space_key, start=0, limit=100)
            all_titles = [p["title"] for p in pages]
            selected_titles = st.multiselect("Select page titles:", options=all_titles)

            for page_title in selected_titles:
                selected_page = next((p for p in pages if p["title"].strip().lower() == page_title.strip().lower()), None)

                if selected_page:
                    page_id = selected_page["id"]
                    html_content = confluence.get_page_by_id(page_id=page_id, expand="body.export_view")["body"]["export_view"]["value"]
                    soup = BeautifulSoup(html_content, "html.parser")
                    base_url = os.getenv("CONFLUENCE_BASE_URL")

                    image_urls = list({
                        base_url + img["src"] if img["src"].startswith("/") else img["src"]
                        for img in soup.find_all("img") if img.get("src")
                    })

                    if not image_urls:
                        st.warning(f"⚠️ No embedded images found in the page '{page_title}'.")
                    else:
                        st.subheader(f"📸 Images in: {page_title}")

                        for idx, url in enumerate(image_urls):
                            summary_key = f"summary_{page_title}_{idx}"
                            ready_key = f"ready_{page_title}_{idx}"
                            image_key = f"image_{page_title}_{idx}"
                            graph_df_key = f"graph_df_{page_title}_{idx}"
                            graph_fig_key = f"graph_fig_{page_title}_{idx}"
                            chart_type_key = f"chart_type_{page_title}_{idx}"

                            try:
                                image_bytes = download_image_bytes(url, auth=(os.getenv('CONFLUENCE_USER_EMAIL'), os.getenv('CONFLUENCE_API_KEY')))
                                st.session_state[image_key] = image_bytes

                                try:
                                    image_pil = Image.open(BytesIO(image_bytes))
                                    st.image(image_pil, use_container_width=True)
                                except UnidentifiedImageError:
                                    st.warning(f"Image {idx} could not be loaded.")
                                    continue

                                if st.button("Summarize", key=f"summarize_{page_title}_{idx}"):
                                    with st.spinner("Generating summary..."):
                                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                                            tmp.write(image_bytes)
                                            tmp.flush()
                                            uploaded = genai.upload_file(
                                                path=tmp.name,
                                                mime_type="image/png",
                                                display_name=f"confluence_image_{page_title}_{idx}.png"
                                            )

                                        prompt = (
                                            "You are analyzing a technical image from a documentation page. "
                                            "If it's a chart or graph, explain what is shown in detail. "
                                            "If it's code, summarize what the code does. "
                                            "Avoid mentioning filenames or metadata. Provide an informative analysis in 2 paragraphs."
                                        )

                                        try:
                                            response = ai_model.generate_content([uploaded, prompt])
                                            st.session_state[summary_key] = response.text.strip()
                                            st.session_state[ready_key] = True           

                                        except Exception as e:
                                            st.error(f"AI failed to generate summary: {e}")

                                if st.session_state.get(ready_key):
                                    anchor_id = f"summary_{page_title}_{idx}"
                                                                                         
                                    st.subheader("Gemini Summary")
                                    summary = st.session_state[summary_key]
                                    st.success(summary)

            
                                    st.markdown("Generate AI response")
                                    ai_response_key = f"ai_response_{page_title}_{idx}"
                                    user_question_key = f"user_question_{page_title}_{idx}"

                                    # Persist user input
                                    if user_question_key not in st.session_state:
                                        st.session_state[user_question_key] = ""

                                    user_question = st.text_input(
                                        "Enter the question:",
                                        key=user_question_key,
                                        value=st.session_state[user_question_key],
                                        placeholder="Ask the Question?"
                                    )

                                    if user_question and st.session_state.get(ai_response_key) is None:
                                        with st.spinner("Generating AI response..."):
                                            try:
                                                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                    tmp_img.write(image_bytes)
                                                    tmp_img.flush()
                                                    uploaded_img = genai.upload_file(
                                                        path=tmp_img.name,
                                                        mime_type="image/png",
                                                        display_name=f"qa_image_{page_title}_{idx}.png"
                                                    )

                                                full_prompt = (
                                                    "You're analyzing a technical image extracted from documentation. "
                                                    "Answer the user's question based on the visual content of the image, "
                                                    "as well as the summary below.\n\n"
                                                    f"Summary:\n{summary}\n\n"
                                                    f"User Question:\n{user_question}"
                                                )

                                                ai_response = ai_model.generate_content([uploaded_img, full_prompt])
                                                st.session_state[ai_response_key] = ai_response.text.strip()
                                            except Exception as e:
                                                st.error(f"AI failed to generate a response: {e}")

                                    # ✅ Display response if available
                                    if st.session_state.get(ai_response_key):
                                        st.info(st.session_state[ai_response_key])

                                    summary = st.session_state[summary_key]

                                    # 🔽 Then the download section
                                    file_name = st.text_input("Enter file name (no extension)", value=f"summary_image_{idx}", key=f"file_name_{page_title}_{idx}")
                                    format_choice = st.selectbox("Select format", ["PDF", "DOCX", "TXT", "Markdown"], key=f"format_{page_title}_{idx}")

                                    content, mime, ext = None, "", ""
                                    if format_choice == "PDF":
                                        content = generate_pdf(image_bytes, summary)
                                        mime = "application/pdf"
                                        ext = "pdf"
                                    elif format_choice == "DOCX":
                                        content = generate_docx(image_bytes, summary)
                                        mime = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                                        ext = "docx"
                                    elif format_choice == "TXT":
                                        content = generate_txt(summary)
                                        mime = "text/plain"
                                        ext = "txt"
                                    elif format_choice == "Markdown":
                                        content = generate_md(image_bytes, summary)
                                        mime = "text/markdown"
                                        ext = "md"

                                    if content:
                                        st.download_button(
                                            label="Download",
                                            data=content,
                                            file_name=f"{file_name}.{ext}",
                                            mime=mime,
                                            key=f"download_button_{page_title}_{idx}"
                                        )

                                        

                                    # ✅ Create Graph button appears only after summary is ready
                                    if st.button("Create Graph", key=f"create_graph_{page_title}_{idx}"):
                                        with st.spinner(" Extracting data from image ..."):
                                            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                                tmp_img.write(image_bytes)
                                                tmp_img.flush()
                                                uploaded_img = genai.upload_file(
                                                    path=tmp_img.name,
                                                    mime_type="image/png",
                                                    display_name=f"chart_image_{page_title}_{idx}.png"
                                                )

                                            graph_prompt = (
                                                "You're looking at a Likert-style bar chart image or table. Extract the full numeric table represented by the chart.\n"
                                                "Return only the raw CSV table: no markdown, no comments, no code blocks.\n"
                                                "The first column must be the response category (e.g., Strongly Agree), followed by columns for group counts (e.g., Students, Lecturers, Staff, Total).\n"
                                                "Ensure all values are numeric and the CSV is properly aligned. Do NOT summarize—just output the table."
                                            )

                                            graph_response = ai_model.generate_content([uploaded_img, graph_prompt])
                                            csv_text = graph_response.text.strip()
                                            

                                            try:
                                                cleaned_csv = clean_ai_csv(csv_text)
                                                df = pd.read_csv(StringIO(cleaned_csv))

                                                for col in df.columns[1:]:
                                                    df[col] = pd.to_numeric(df[col], errors='coerce')

                                                df.dropna(subset=df.columns[1:], how='all', inplace=True)

                                                if df.empty:
                                                    raise ValueError("Extracted DataFrame is empty after cleaning.")

                                                st.session_state[graph_df_key] = df
                                                st.subheader("Table from Image")
                                                st.dataframe(df)
                                                

                                            except Exception as e:
                                                st.error(f"⚠️ Failed to parse chart data. Reason: {e}")
                                                st.text_area("🧾 AI Output ", csv_text, height=200)

                                if graph_df_key in st.session_state:
                                    df = st.session_state[graph_df_key]
                                    chart_type = st.selectbox("📈 Choose Chart Type", ["Grouped Bar", "Stacked Bar", "Line", "Pie"], key=chart_type_key)

                                    if chart_type == "Grouped Bar":
                                        fig = plot_grouped_bar(df)
                                    elif chart_type == "Stacked Bar":
                                        fig = plot_stacked_bar(df)
                                    elif chart_type == "Line":
                                        fig = plot_line(df)
                                    else:
                                        fig = plot_pie(df)

                                    st.session_state[graph_fig_key] = fig
                                    st.pyplot(fig)
                                    
                                    

                                    download_format = st.selectbox(
                                        "Download Format",
                                        ["PNG", "JPG", "SVG", "PDF", "DOCX", "PPTX"],
                                        key=f"download_fmt_{page_title}_{idx}"
                                    )

                                    chart_file_name = st.text_input("Enter chart file name", f"chart_{page_title}_{idx}", key=f"chart_file_name_{page_title}_{idx}")

                                    # Default values
                                    fmt_ext, mime, download_data = None, None, None

                                    if download_format == "PNG":
                                        fmt_ext, mime = "png", "image/png"
                                        download_data = get_image_bytes(fig, "png")
                                    elif download_format == "JPG":
                                        fmt_ext, mime = "jpg", "image/jpeg"
                                        download_data = get_image_bytes(fig, "jpg")
                                    elif download_format == "SVG":
                                        fmt_ext, mime = "svg", "image/svg+xml"
                                        download_data = get_image_bytes(fig, "svg")
                                    elif download_format == "PDF":
                                        fmt_ext, mime = "pdf", "application/pdf"
                                        download_data = get_image_bytes(fig, "pdf")
                                    elif download_format == "DOCX":
                                        from docx import Document
                                        from docx.shared import Inches

                                        output = BytesIO()
                                        doc = Document()
                                        image_stream = get_image_bytes(fig, "png")
                                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                            tmp_img.write(image_stream.read())
                                            tmp_img.flush()
                                            doc.add_picture(tmp_img.name, width=Inches(6))
                                        doc.save(output)
                                        output.seek(0)
                                        fmt_ext, mime, download_data = "docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", output
                                    elif download_format == "PPTX":
                                        from pptx import Presentation
                                        from pptx.util import Inches

                                        prs = Presentation()
                                        slide_layout = prs.slide_layouts[5]  # blank layout
                                        slide = prs.slides.add_slide(slide_layout)

                                        image_stream = get_image_bytes(fig, "png")
                                        with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp_img:
                                            tmp_img.write(image_stream.read())
                                            tmp_img.flush()
                                            slide.shapes.add_picture(tmp_img.name, Inches(1), Inches(1), Inches(6), Inches(4.5))

                                        output = BytesIO()
                                        prs.save(output)
                                        output.seek(0)
                                        fmt_ext, mime, download_data = "pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", output

                                    if download_data:
                                        st.download_button(
                                            label="Download Graph",
                                            data=download_data,
                                            file_name=f"{chart_file_name}.{fmt_ext}",
                                            mime=mime,
                                            key=f"download_chart_{page_title}_{idx}"
                                        )

                                    
                            except Exception as e:
                                st.error(f"⚠️ Error processing image {idx}: {str(e)}")
        except Exception as e:
            st.error(f"❌ Error fetching pages or processing images: {str(e)}")
else:
    st.error("❌ Failed to connect to Confluence.")
